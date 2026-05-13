"""PPTX 内容提取器"""

from pathlib import Path


class PptxExtractor:
    """PowerPoint 文档内容提取"""

    SUPPORTED_EXTENSIONS = {".pptx", ".ppt"}

    def extract_text(self, file_path: str | Path) -> str:
        """提取 PPTX 中的文本内容"""
        try:
            from pptx import Presentation
        except ImportError:
            return ""

        try:
            prs = Presentation(str(file_path))
            parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_text.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                slide_text.append(" | ".join(cells))
                if slide_text:
                    parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
            return "\n\n".join(parts)
        except Exception:
            return ""

    def extract_metadata(self, file_path: str | Path) -> dict:
        """提取 PPTX 元数据"""
        try:
            from pptx import Presentation
        except ImportError:
            return {}

        try:
            prs = Presentation(str(file_path))
            meta = {
                "slide_count": len(prs.slides),
            }
            if prs.core_properties:
                meta["title"] = prs.core_properties.title or ""
                meta["author"] = prs.core_properties.author or ""
            return meta
        except Exception:
            return {}
